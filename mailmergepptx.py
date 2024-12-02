#Dependencies
#pip install python-pptx openpyxl comtypes six PyPDF2

import sys
import openpyxl
from pptx import Presentation
import copy
import os
import six
import comtypes.client
from PyPDF2 import PdfMerger


def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

def merge_ppt(excel_filename, ppt_filename):
    # Load the Excel file
    workbook = openpyxl.load_workbook(excel_filename)
    sheet = workbook.active

    # Get column headers from the first row
    headers = [cell.value for cell in sheet[1]]

    # Load the PowerPoint template
    prs = Presentation(ppt_filename)

    output_filename = f"merged_{ppt_filename}"
    print(f"Generating {output_filename}...")

    # Iterate over each row in the Excel sheet starting from the second row
    for row_number, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), start=2):
        replacements = {f"{headers[i]}": row[i] for i in range(len(headers))}

        # Get the slide number corresponding to the row number
        slide_num = row_number - 2  # Assuming the first row (min_row=2) corresponds to the first slide (index 0)
        
        # Check if the slide number is within the range of slides in the presentation
        if slide_num >= len(prs.slides):
            continue  # Skip if slide number is out of range
        
        slide = prs.slides[slide_num]

        # Replace text in the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                replace_text_in_shape(shape, replacements)

        # Save the modified presentation

        prs.save(output_filename)
        print(f"Updated row {row_number:03} to slide {row_number - 1:03}")
    print(f"Generated {output_filename}")

def generate_split_ppt(excel_filename, ppt_filename):
    # Load the Excel file
    workbook = openpyxl.load_workbook(excel_filename)
    sheet = workbook.active

    # Get column headers from the first rows
    headers = [cell.value for cell in sheet[1]]

    # Iterate over each row in the Excel sheet starting from the second row
    for row_number, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), start=2):
        replacements = {f"{headers[i]}": row[i] for i in range(len(headers))}

        # Load the PowerPoint template
        prs = Presentation(ppt_filename)

        # Replace text in each slide
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    replace_text_in_shape(shape, replacements)

        # Save the modified presentation
        output_filename = f"Row_{row_number:03}_slide.pptx"
        prs.save(output_filename)
        print(f"Generated {output_filename}")
    convert_all_pptx_to_pdf()
    merge_pdfs()

def convert_all_pptx_to_pdf():
    try:
        folder_path = os.getcwd()
        for filename in os.listdir(folder_path):
            if filename.endswith("_slide.pptx"):
                pdf_path = filename[:filename.rfind('.')]
                convert_pptx_to_pdf(filename, f'{pdf_path}.pdf', True)
        #Close powerpoint application
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Quit()
    except Exception as e:
        print(f"Failed to generate pdf. Error: {e}")

def merge_pdfs():
    try:
        folder_path = os.getcwd()
        output_filename = "merged.pdf"
        merger = PdfMerger()

        for filename in os.listdir(folder_path):
            if filename.endswith(".pdf") and filename.find('merged.pdf') < 0:
                print(f"Merging {filename}")
                pdf_path = os.path.join(folder_path, filename)
                merger.append(pdf_path)

        with open(output_filename, 'wb') as output_file:
            merger.write(output_file)
            merger.close()
        print(f"Merged all PDFs into {output_filename}")
    except Exception as e:
        print(f"Failed to merge pdfs in the folder to merged.pdf. Error: {e}")

def convert_pptx_to_pdf(pptx_path, pdf_path, inLoop=False): 
    try:
        print(f"Generating {pptx_path} to {pdf_path}. Opening slide in Powerpoint...")
        # Initialize PowerPoint
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1

        # Open the presentation
        presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path))

        # Save the presentation as PDF
        print(f"Saving {pdf_path}")
        presentation.SaveAs(os.path.abspath(pdf_path), 32)  # 32 is the constant for PDF format
        presentation.Close()
        if not inLoop:
            powerpoint.Quit()
        print(f"Converted {pptx_path} to {pdf_path}")
    except Exception as e:
        print(f"Failed to convert {pptx_path} to PDF. Error: {e}")
        powerpoint.Quit()

if __name__ == "__main__":
    print("mailmergepptx Version: 0.1")
    print("Author: bopaiah@hotmail.com")

    argslen = len(sys.argv)
    if "/p" in sys.argv and argslen == 3:
        print("Converting pptx to pdf")
        pptx_path = sys.argv[2]
        if pptx_path.upper().endswith('PPTX'):
            pdf_path = pptx_path[:pptx_path.rfind('.')]
            convert_pptx_to_pdf(pptx_path, f'{pdf_path}.pdf')
        else:
            print('Note: Only PPTX file can be converted to PDF')
    elif "/a" in sys.argv:
        convert_all_pptx_to_pdf()
    elif "/m" in sys.argv:
        merge_pdfs()
    elif "/s" in sys.argv and argslen == 4:
        excel_filename = sys.argv[2]
        ppt_filename = sys.argv[3]
        generate_split_ppt(excel_filename, ppt_filename)
    elif argslen == 3:
        excel_filename = sys.argv[1]
        ppt_filename = sys.argv[2]
        merge_ppt(excel_filename, ppt_filename)
    else:
        print("Usage: python mailmergepptx <excel_filename> <ppt_filename>")
        print("       Use the same name as the column headers in Excel for replacing the text in the slide.\n")
        print("       Create a pptx file with as many slides as the number of rows in excel. One slide for each row.")
        print("       The output will be merged_filename.pptx with data from excel row replaced into the corresponding slide number.")
        print("       Program works only on Windows with MS PowerPoint installation is required for the exporting pptx to pdf.")
        print("       Example: If the column header is 'title', the application will replace all the text 'title'")
        print("                throughout the slide with the content in the excel rows.\n")
        print("       /p <pptx filename> for only converting pptx to pdf. The PDF file is generated in same folder as PPTX.")
        print("       /a  with no other arguments for converting all the generated *_slide.pptx to pdf.")
        print("       /m  with no other arguments for merging all pdf files in the current folder to merged.pdf.")
        print("       /s  Generates Single ppt for each row from excel and converts to PDF.")

# End of program
